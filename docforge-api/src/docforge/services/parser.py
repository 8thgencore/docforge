from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path


@dataclass(slots=True)
class ParsedDocument:
    text: str
    language: str | None
    metadata: dict


class ParserError(RuntimeError):
    pass


def parse_document(file_path: Path) -> ParsedDocument:
    suffix = file_path.suffix.lower()
    if suffix in {".txt", ".md", ".csv", ".json", ".yaml", ".yml"}:
        return ParsedDocument(text=file_path.read_text(errors="ignore"), language=None, metadata={})
    if suffix == ".pdf":
        return _parse_pdf(file_path)
    if suffix == ".docx":
        return _parse_docx(file_path)
    if suffix == ".xlsx":
        return _parse_xlsx(file_path)
    if suffix == ".pptx":
        return _parse_pptx(file_path)

    try:
        text = file_path.read_text(errors="ignore")
        return ParsedDocument(text=text, language=None, metadata={"parser": "fallback_text"})
    except Exception as exc:  # pragma: no cover - hard fallback
        raise ParserError(f"unsupported file type: {file_path.name}") from exc


def _parse_pdf(file_path: Path) -> ParsedDocument:
    try:
        import fitz
    except Exception as exc:  # pragma: no cover - dependency issue
        raise ParserError("PyMuPDF is required for PDF parsing") from exc

    with fitz.open(file_path) as doc:
        text = "\n".join(page.get_text() for page in doc)
    return ParsedDocument(text=text, language=None, metadata={"parser": "pymupdf"})


def _parse_docx(file_path: Path) -> ParsedDocument:
    try:
        from docx import Document
    except Exception as exc:  # pragma: no cover
        raise ParserError("python-docx is required for DOCX parsing") from exc

    document = Document(file_path)
    text = "\n".join(paragraph.text for paragraph in document.paragraphs)
    return ParsedDocument(text=text, language=None, metadata={"parser": "python-docx"})


def _parse_xlsx(file_path: Path) -> ParsedDocument:
    try:
        from openpyxl import load_workbook
    except Exception as exc:  # pragma: no cover
        raise ParserError("openpyxl is required for XLSX parsing") from exc

    workbook = load_workbook(file_path, data_only=True)
    values: list[str] = []
    for sheet in workbook.worksheets:
        values.append(f"# {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(cell) for cell in row if cell is not None]
            if cells:
                values.append(" | ".join(cells))
    return ParsedDocument(text="\n".join(values), language=None, metadata={"parser": "openpyxl"})


def _parse_pptx(file_path: Path) -> ParsedDocument:
    try:
        from pptx import Presentation
    except Exception as exc:  # pragma: no cover
        raise ParserError("python-pptx is required for PPTX parsing") from exc

    presentation = Presentation(file_path)
    lines: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        lines.append(f"# Slide {slide_index}")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    lines.append(text)
    return ParsedDocument(text="\n".join(lines), language=None, metadata={"parser": "python-pptx"})
